from pptx import Presentation
import sys

prs = Presentation('/Users/parthgupta/Downloads/FIN Project/Finspark_Hackathon_Template.pptx')
for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_format = shape.placeholder_format
            text = shape.text if hasattr(shape, 'text') else ''
            print(f"  Placeholder {ph_format.idx}: type={ph_format.type}, name={shape.name}, text={repr(text)}")
        elif hasattr(shape, 'text'):
            print(f"  Shape (no PH): name={shape.name}, text={repr(shape.text)}")
