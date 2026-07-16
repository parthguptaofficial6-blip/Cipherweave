import collections
import collections.abc
from pptx import Presentation

# Fix for python-pptx with python 3.10+
import pptx
pptx.collections = collections

def replace_text_in_shape(shape, new_text):
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = new_text

prs = Presentation('/Users/parthgupta/Downloads/FIN Project/Finspark_Hackathon_Template.pptx')

slides_content = {
    0: {
        "Your Team Name": "Team Name: Team CipherWeave\n\nTeam Bio: Innovators in cyber-financial risk correlation.\n\nDate: July 2026"
    },
    1: {
        "problem statement": "Problem: Financial institutions struggle to connect cyber threat intelligence (e.g., unusual logins, weak cipher suites) with financial fraud. Siloed data leads to delayed detection of coordinated attacks.\n\nWhy we chose it: By correlating these streams using AI, we can drastically reduce false positives and catch sophisticated multi-vector attacks (like Harvest Now, Decrypt Later)."
    },
    2: {
        "assumptions": "• Node.js & npm (Frontend: React + Vite)\n• Python 3.8+ (Backend: FastAPI, pandas, scikit-learn, networkx)\n• SQLite (Local database)"
    },
    3: {
        "technologies": "Frontend: React, Vite, Tailwind CSS, Lucide React\nBackend: Python, FastAPI, SQLAlchemy, Uvicorn\nData & AI: NetworkX (Graph Correlation), Scikit-Learn (Isolation Forest), Pandas, NumPy"
    },
    4: {
        "supporting documents": "Explaining the logic flow:\n1. Cyber and Transaction events are simulated.\n2. Correlated using Graph Networks (Entity ID).\n3. Feature engineered (Temporal proximity, Geo anomaly, Amount Z-score, Quantum risk).\n4. Scored via Isolation Forest.\n5. Dashboard display."
    },
    5: {
        "technologies": "Differentiators: We combine traditional financial monitoring with cyber indicators (like TLS cipher suite monitoring for Quantum threats).\n\nAdoption: Easy to integrate via APIs. Plug-and-play dashboards for risk analysts."
    },
    6: {
        "GitHub link": "GitHub Link: (Added manually by team)\n\nNote: Please insert architecture diagrams or screenshots here."
    },
    7: {
        "future potential": "• Reduces financial losses by detecting coordinated multi-channel attacks early.\n• Prepares institutions for emerging threats like Quantum computing (HNDL attacks).\n• Can be expanded to incorporate more data sources (e.g., social engineering indicators)."
    },
    8: {
        "innovative": "• Uses Graph-based Entity Resolution to link siloed cyber events and financial transactions.\n• Provides ML-backed explainability (Explainable AI) so analysts know *why* a case was flagged (e.g., '+0.45 Temporal Proximity')."
    },
    9: {
        "how users will interact": "• The React-based dashboard is clean, fast, and interactive.\n• Analysts can quickly expand high-risk cases, review the exact correlated events, and take immediate action (Escalate/Dismiss/False Positive)."
    },
    10: {
        "how the solution can scale": "• Backend: FastAPI handles high-throughput asynchronous requests.\n• AI Engine: Isolation forest and NetworkX graphs can be distributed using tools like Spark/Dask for enterprise-scale transaction volumes.\n• Database: Can easily migrate from SQLite to PostgreSQL/CockroachDB."
    },
    11: {
        "how practical the solution is": "• Backend uses standard Python virtual environments and requirements.txt.\n• Frontend uses Vite for lightning-fast builds.\n• Architecture is decoupled (REST API), making containerization (Docker) and Kubernetes orchestration straightforward."
    },
    12: {
        "cybersecurity risks": "• Includes JWT-based authentication for securing the API.\n• Recommends strong password hashing (bcrypt).\n• Implements CORS middleware.\n• The solution actively monitors for deprecated cipher suites (TLS_RSA_WITH_AES_128_CBC_SHA) which are vulnerable to future quantum attacks."
    },
    13: {
        "architecture diagram": "[Please insert your architecture diagram here]\n\nFrontend (React) <-> REST API (FastAPI) <-> Database (SQLite) & Correlation Engine (Scikit-Learn/NetworkX)."
    },
    14: {
        "prototype screenshots": "[Please insert your solution screenshots and demo video link here]\n\nDashboard displays synthetic risk-scored cases."
    },
    15: {
        "team member names": "Team CipherWeave"
    }
}

for i, slide in enumerate(prs.slides):
    if i in slides_content:
        content_map = slides_content[i]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape_text = shape.text
                for key, new_text in content_map.items():
                    if key in shape_text:
                        replace_text_in_shape(shape, new_text)

prs.save('/Users/parthgupta/Downloads/FIN Project/CipherWeave_Submission.pptx')
print("Successfully generated CipherWeave_Submission.pptx")
